"""Phase O2 — Office & EML Reconnaissance tests.

Tests cover:
- files/office_docx.py: DOCX extraction
- files/office_xlsx.py: XLSX extraction
- files/office_pptx.py: PPTX extraction
- files/email_eml.py: EML extraction
- tools/file_inspect.py: routing for new types + legacy degradation
"""

import json
import os
from unittest.mock import AsyncMock

import pytest

from tools.base import ToolContext


# ── Helpers ──────────────────────────────────────────────────────────────────


def _make_ctx(session_dir: str) -> ToolContext:
    return ToolContext(
        user_id="test-user",
        session_id="test-session",
        user_root=session_dir,
        session_dir=session_dir,
        workspace_dir=session_dir,
        venv_bin="",
        publish=AsyncMock(),
    )


def _create_docx(path: str, paragraphs: list[str] | None = None, tables: int = 0) -> str:
    from docx import Document

    doc = Document()
    doc.add_heading("Test Document", level=1)
    doc.add_heading("Introduction", level=2)
    for text in (paragraphs or ["This is a test paragraph with enough content."]):
        doc.add_paragraph(text)
    for _ in range(tables):
        doc.add_table(rows=2, cols=3)
    doc.save(path)
    return path


def _create_xlsx(path: str, sheets: dict[str, list[list]] | None = None) -> str:
    from openpyxl import Workbook

    wb = Workbook()
    data = sheets or {"Sheet1": [["Name", "Age", "City"], ["Alice", 30, "NYC"], ["Bob", 25, "LA"]]}
    first = True
    for name, rows in data.items():
        if first:
            ws = wb.active
            ws.title = name
            first = False
        else:
            ws = wb.create_sheet(name)
        for row in rows:
            ws.append(row)
    wb.save(path)
    return path


def _create_pptx(path: str, slide_count: int = 3) -> str:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    for i in range(slide_count):
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # title + content
        slide.shapes.title.text = f"Slide {i + 1} Title"
        slide.placeholders[1].text = f"Content for slide {i + 1}. This is sample text."
    prs.save(path)
    return path


def _create_eml(path: str, subject: str = "Test Email", body: str = "Hello, this is a test email.", attachments: list[str] | None = None) -> str:
    from email.message import EmailMessage

    msg = EmailMessage()
    msg["Subject"] = subject
    msg["From"] = "sender@example.com"
    msg["To"] = "recipient@example.com"
    msg["Date"] = "Mon, 31 Mar 2026 10:00:00 +0800"
    msg.set_content(body)

    for att_name in (attachments or []):
        msg.add_attachment(b"dummy content", maintype="application", subtype="octet-stream", filename=att_name)

    with open(path, "wb") as f:
        f.write(msg.as_bytes())
    return path


# ── Test: DOCX extraction ───────────────────────────────────────────────────


class TestDocxExtract:

    def test_basic_structure(self, tmp_path):
        from files.office_docx import extract

        path = _create_docx(str(tmp_path / "test.docx"))
        result = extract(path)

        assert result["kind"] == "office"
        assert result["office_kind"] == "docx"
        assert result["size_bytes"] > 0
        assert result["paragraph_count"] >= 1
        assert result["heading_count"] >= 1
        assert isinstance(result["headings"], list)
        assert isinstance(result["text_sample"], str)
        assert isinstance(result["metadata"], dict)

    def test_headings_extracted(self, tmp_path):
        from files.office_docx import extract

        path = _create_docx(str(tmp_path / "test.docx"))
        result = extract(path)

        assert "Test Document" in result["headings"]
        assert "Introduction" in result["headings"]

    def test_tables_counted(self, tmp_path):
        from files.office_docx import extract

        path = _create_docx(str(tmp_path / "test.docx"), tables=2)
        result = extract(path)

        assert result["table_count"] == 2

    def test_file_not_found(self, tmp_path):
        from files.office_docx import extract

        with pytest.raises(FileNotFoundError):
            extract(str(tmp_path / "nope.docx"))

    def test_invalid_docx(self, tmp_path):
        from files.office_docx import extract

        bad = str(tmp_path / "bad.docx")
        with open(bad, "w") as f:
            f.write("not a docx")
        with pytest.raises(ValueError, match="Cannot read DOCX"):
            extract(bad)


# ── Test: XLSX extraction ───────────────────────────────────────────────────


class TestXlsxExtract:

    def test_basic_structure(self, tmp_path):
        from files.office_xlsx import extract

        path = _create_xlsx(str(tmp_path / "test.xlsx"))
        result = extract(path)

        assert result["kind"] == "office"
        assert result["office_kind"] == "xlsx"
        assert result["sheet_count"] == 1
        assert result["sheet_names"] == ["Sheet1"]
        assert len(result["sheets"]) == 1

    def test_sheet_details(self, tmp_path):
        from files.office_xlsx import extract

        path = _create_xlsx(str(tmp_path / "test.xlsx"))
        result = extract(path)

        sheet = result["sheets"][0]
        assert sheet["name"] == "Sheet1"
        assert sheet["header_row"] == ["Name", "Age", "City"]
        assert len(sheet["sample_rows"]) == 2
        assert sheet["sample_rows"][0] == ["Alice", "30", "NYC"]

    def test_multiple_sheets(self, tmp_path):
        from files.office_xlsx import extract

        path = _create_xlsx(str(tmp_path / "multi.xlsx"), sheets={
            "Sales": [["Date", "Amount"], ["2026-01", "100"]],
            "Summary": [["Total"], ["100"]],
        })
        result = extract(path)

        assert result["sheet_count"] == 2
        assert result["sheet_names"] == ["Sales", "Summary"]

    def test_file_not_found(self, tmp_path):
        from files.office_xlsx import extract

        with pytest.raises(FileNotFoundError):
            extract(str(tmp_path / "nope.xlsx"))

    def test_invalid_xlsx(self, tmp_path):
        from files.office_xlsx import extract

        bad = str(tmp_path / "bad.xlsx")
        with open(bad, "w") as f:
            f.write("not an xlsx")
        with pytest.raises(ValueError, match="Cannot read XLSX"):
            extract(bad)


# ── Test: PPTX extraction ──────────────────────────────────────────────────


class TestPptxExtract:

    def test_basic_structure(self, tmp_path):
        from files.office_pptx import extract

        path = _create_pptx(str(tmp_path / "test.pptx"))
        result = extract(path)

        assert result["kind"] == "office"
        assert result["office_kind"] == "pptx"
        assert result["slide_count"] == 3
        assert len(result["slides"]) == 3
        assert isinstance(result["metadata"], dict)

    def test_slide_details(self, tmp_path):
        from files.office_pptx import extract

        path = _create_pptx(str(tmp_path / "test.pptx"), slide_count=2)
        result = extract(path)

        slide = result["slides"][0]
        assert slide["number"] == 1
        assert "Slide 1 Title" in slide["title"]
        assert isinstance(slide["text_preview"], str)
        assert isinstance(slide["has_notes"], bool)

    def test_file_not_found(self, tmp_path):
        from files.office_pptx import extract

        with pytest.raises(FileNotFoundError):
            extract(str(tmp_path / "nope.pptx"))

    def test_invalid_pptx(self, tmp_path):
        from files.office_pptx import extract

        bad = str(tmp_path / "bad.pptx")
        with open(bad, "w") as f:
            f.write("not a pptx")
        with pytest.raises(ValueError, match="Cannot read PPTX"):
            extract(bad)


# ── Test: EML extraction ───────────────────────────────────────────────────


class TestEmlExtract:

    def test_basic_structure(self, tmp_path):
        from files.email_eml import extract

        path = _create_eml(str(tmp_path / "test.eml"))
        result = extract(path)

        assert result["kind"] == "email"
        assert result["email_kind"] == "eml"
        assert result["subject"] == "Test Email"
        assert "sender@example.com" in result["from_addr"]
        assert "recipient@example.com" in result["to_addr"]
        assert result["date"] != ""
        assert "test email" in result["body_preview"].lower()

    def test_attachments(self, tmp_path):
        from files.email_eml import extract

        path = _create_eml(str(tmp_path / "att.eml"), attachments=["report.pdf", "data.xlsx"])
        result = extract(path)

        assert result["attachment_count"] == 2
        filenames = [a["filename"] for a in result["attachments"]]
        assert "report.pdf" in filenames
        assert "data.xlsx" in filenames

    def test_no_attachments(self, tmp_path):
        from files.email_eml import extract

        path = _create_eml(str(tmp_path / "plain.eml"))
        result = extract(path)

        assert result["attachment_count"] == 0
        assert result["attachments"] == []

    def test_file_not_found(self, tmp_path):
        from files.email_eml import extract

        with pytest.raises(FileNotFoundError):
            extract(str(tmp_path / "nope.eml"))


# ── Test: file_inspect routing for O2 types ─────────────────────────────────


class TestFileInspectO2:

    @pytest.mark.asyncio
    async def test_inspect_docx(self, tmp_path):
        from tools.file_inspect import FileInspectTool

        _create_docx(str(tmp_path / "doc.docx"))
        tool = FileInspectTool()
        ctx = _make_ctx(str(tmp_path))

        result = await tool.execute(ctx, path="doc.docx")
        assert result["is_error"] is False
        data = json.loads(result["output"])
        assert data["kind"] == "office"
        assert data["office_kind"] == "docx"
        assert data["understanding_available"] is True

    @pytest.mark.asyncio
    async def test_inspect_xlsx(self, tmp_path):
        from tools.file_inspect import FileInspectTool

        _create_xlsx(str(tmp_path / "data.xlsx"))
        tool = FileInspectTool()
        ctx = _make_ctx(str(tmp_path))

        result = await tool.execute(ctx, path="data.xlsx")
        assert result["is_error"] is False
        data = json.loads(result["output"])
        assert data["kind"] == "office"
        assert data["office_kind"] == "xlsx"
        assert data["understanding_available"] is True

    @pytest.mark.asyncio
    async def test_inspect_pptx(self, tmp_path):
        from tools.file_inspect import FileInspectTool

        _create_pptx(str(tmp_path / "slides.pptx"))
        tool = FileInspectTool()
        ctx = _make_ctx(str(tmp_path))

        result = await tool.execute(ctx, path="slides.pptx")
        assert result["is_error"] is False
        data = json.loads(result["output"])
        assert data["kind"] == "office"
        assert data["office_kind"] == "pptx"
        assert data["understanding_available"] is True

    @pytest.mark.asyncio
    async def test_inspect_eml(self, tmp_path):
        from tools.file_inspect import FileInspectTool

        _create_eml(str(tmp_path / "mail.eml"))
        tool = FileInspectTool()
        ctx = _make_ctx(str(tmp_path))

        result = await tool.execute(ctx, path="mail.eml")
        assert result["is_error"] is False
        data = json.loads(result["output"])
        assert data["kind"] == "email"
        assert data["email_kind"] == "eml"
        assert data["understanding_available"] is True


# ── Test: Legacy degradation ────────────────────────────────────────────────


class TestLegacyDegradation:

    @pytest.mark.asyncio
    async def test_doc_degradation(self, tmp_path):
        from tools.file_inspect import FileInspectTool

        # Create a dummy .doc file
        doc_path = str(tmp_path / "old.doc")
        with open(doc_path, "wb") as f:
            f.write(b"\xd0\xcf\x11\xe0")  # OLE magic bytes
        tool = FileInspectTool()
        ctx = _make_ctx(str(tmp_path))

        result = await tool.execute(ctx, path="old.doc")
        assert result["is_error"] is False
        data = json.loads(result["output"])
        assert data["kind"] == "office"
        assert data["office_kind"] == "doc"
        assert data["understanding_available"] is False
        assert "docx" in data["message"].lower()

    @pytest.mark.asyncio
    async def test_xls_degradation(self, tmp_path):
        from tools.file_inspect import FileInspectTool

        xls_path = str(tmp_path / "old.xls")
        with open(xls_path, "wb") as f:
            f.write(b"\xd0\xcf\x11\xe0")
        tool = FileInspectTool()
        ctx = _make_ctx(str(tmp_path))

        result = await tool.execute(ctx, path="old.xls")
        assert result["is_error"] is False
        data = json.loads(result["output"])
        assert data["kind"] == "office"
        assert data["office_kind"] == "xls"
        assert data["understanding_available"] is False

    @pytest.mark.asyncio
    async def test_ppt_degradation(self, tmp_path):
        from tools.file_inspect import FileInspectTool

        ppt_path = str(tmp_path / "old.ppt")
        with open(ppt_path, "wb") as f:
            f.write(b"\xd0\xcf\x11\xe0")
        tool = FileInspectTool()
        ctx = _make_ctx(str(tmp_path))

        result = await tool.execute(ctx, path="old.ppt")
        assert result["is_error"] is False
        data = json.loads(result["output"])
        assert data["understanding_available"] is False

    @pytest.mark.asyncio
    async def test_msg_degradation(self, tmp_path):
        from tools.file_inspect import FileInspectTool

        msg_path = str(tmp_path / "old.msg")
        with open(msg_path, "wb") as f:
            f.write(b"\xd0\xcf\x11\xe0")
        tool = FileInspectTool()
        ctx = _make_ctx(str(tmp_path))

        result = await tool.execute(ctx, path="old.msg")
        assert result["is_error"] is False
        data = json.loads(result["output"])
        assert data["kind"] == "email"
        assert data["email_kind"] == "msg"
        assert data["understanding_available"] is False
        assert "eml" in data["message"].lower()


# ── Test: Registry still correct ────────────────────────────────────────────


class TestRegistryO2:

    def test_tool_count_is_11(self):
        from tools.registry import get_registry

        registry = get_registry()
        assert len(registry.tools) == 22

    def test_file_inspect_description_mentions_office(self):
        from tools.registry import get_registry

        registry = get_registry()
        tool = registry.get("file_inspect")
        assert "Office" in tool.description or "DOCX" in tool.description


if __name__ == "__main__":
    pytest.main([__file__, "-v"])
