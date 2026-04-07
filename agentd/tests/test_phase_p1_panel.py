"""Phase P1 — Panel Workbench backend tests.

Tests cover:
- workspace/inspect endpoint (file extraction via REST, independent of agent loop)
- Panel Content Protocol schemas
- _extract_file_info dispatch
"""

import json
import os

import pytest

from workspace.manager import ensure_user_root, get_session_dir
from workspace.schemas import (
    HtmlSandboxContent,
    PanelContent,
    PanelContentType,
    PanelType,
    StructuredContent,
    StructuredWidget,
)


# ── Fixtures ──────────────────────────────────────────────────────────────────


@pytest.fixture
def user_root(tmp_path):
    root = os.path.join(str(tmp_path), "test-user")
    ensure_user_root(root)
    return root


@pytest.fixture
def session_dir(user_root):
    sd = get_session_dir(user_root, "test-session")
    os.makedirs(os.path.join(sd, ".agentd"), exist_ok=True)
    return sd


# ── _extract_file_info dispatch ──────────────────────────────────────────


class TestExtractFileInfo:
    @pytest.mark.asyncio
    async def test_pdf_extraction(self, session_dir):
        """PDF files dispatch to files.pdf.extract."""
        from pypdf import PdfWriter
        pdf_path = os.path.join(session_dir, "test.pdf")
        writer = PdfWriter()
        writer.add_blank_page(width=612, height=792)
        with open(pdf_path, "wb") as f:
            writer.write(f)

        from workspace.router import _extract_file_info
        result = await _extract_file_info(pdf_path, ".pdf")

        assert result["kind"] == "pdf"
        assert "page_count" in result
        assert result["page_count"] == 1

    @pytest.mark.asyncio
    async def test_docx_extraction(self, session_dir):
        """DOCX files dispatch to files.office_docx.extract."""
        from docx import Document
        docx_path = os.path.join(session_dir, "test.docx")
        doc = Document()
        doc.add_heading("Test Title", level=1)
        doc.add_paragraph("Test paragraph content.")
        doc.save(docx_path)

        from workspace.router import _extract_file_info
        result = await _extract_file_info(docx_path, ".docx")

        assert result["kind"] == "office"
        assert "paragraph_count" in result

    @pytest.mark.asyncio
    async def test_xlsx_extraction(self, session_dir):
        """XLSX files dispatch to files.office_xlsx.extract."""
        from openpyxl import Workbook
        xlsx_path = os.path.join(session_dir, "test.xlsx")
        wb = Workbook()
        ws = wb.active
        ws.append(["Name", "Age", "City"])
        ws.append(["Alice", 30, "Beijing"])
        wb.save(xlsx_path)

        from workspace.router import _extract_file_info
        result = await _extract_file_info(xlsx_path, ".xlsx")

        assert result["kind"] == "office"
        assert "sheets" in result

    @pytest.mark.asyncio
    async def test_pptx_extraction(self, session_dir):
        """PPTX files dispatch to files.office_pptx.extract."""
        from pptx import Presentation
        pptx_path = os.path.join(session_dir, "test.pptx")
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Slide Title"
        prs.save(pptx_path)

        from workspace.router import _extract_file_info
        result = await _extract_file_info(pptx_path, ".pptx")

        assert result["kind"] == "office"
        assert "slide_count" in result

    @pytest.mark.asyncio
    async def test_eml_extraction(self, session_dir):
        """EML files dispatch to files.email_eml.extract."""
        eml_path = os.path.join(session_dir, "test.eml")
        eml_content = (
            "From: sender@example.com\r\n"
            "To: receiver@example.com\r\n"
            "Subject: Test Email\r\n"
            "Date: Mon, 31 Mar 2026 10:00:00 +0800\r\n"
            "\r\n"
            "This is the email body.\r\n"
        )
        with open(eml_path, "w") as f:
            f.write(eml_content)

        from workspace.router import _extract_file_info
        result = await _extract_file_info(eml_path, ".eml")

        assert result["kind"] == "email"
        assert result["subject"] == "Test Email"

    @pytest.mark.asyncio
    async def test_image_extraction(self, session_dir):
        """Image files dispatch to files.image.extract_metadata."""
        from PIL import Image
        img_path = os.path.join(session_dir, "test.png")
        img = Image.new("RGB", (200, 100), color=(255, 0, 0))
        img.save(img_path)

        from workspace.router import _extract_file_info
        result = await _extract_file_info(img_path, ".png")

        assert result["kind"] == "image"
        assert result["width"] == 200
        assert result["height"] == 100

    @pytest.mark.asyncio
    async def test_file_not_found(self, session_dir):
        from workspace.router import _extract_file_info
        result = await _extract_file_info("/nonexistent/file.pdf", ".pdf")
        assert result["kind"] == "error"

    @pytest.mark.asyncio
    async def test_unknown_extension(self, session_dir):
        from workspace.router import _extract_file_info
        result = await _extract_file_info("/any/file.xyz", ".xyz")
        assert result["kind"] == "unknown"


# ── Inspectable extensions ───────────────────────────────────────────────


class TestInspectableExtensions:
    def test_supported_extensions(self):
        from workspace.router import _INSPECTABLE_EXTENSIONS
        assert ".pdf" in _INSPECTABLE_EXTENSIONS
        assert ".docx" in _INSPECTABLE_EXTENSIONS
        assert ".xlsx" in _INSPECTABLE_EXTENSIONS
        assert ".pptx" in _INSPECTABLE_EXTENSIONS
        assert ".eml" in _INSPECTABLE_EXTENSIONS
        assert ".png" in _INSPECTABLE_EXTENSIONS
        assert ".jpg" in _INSPECTABLE_EXTENSIONS

    def test_unsupported_not_in_set(self):
        from workspace.router import _INSPECTABLE_EXTENSIONS
        assert ".txt" not in _INSPECTABLE_EXTENSIONS
        assert ".py" not in _INSPECTABLE_EXTENSIONS
        assert ".doc" not in _INSPECTABLE_EXTENSIONS  # legacy, not inspectable


# ── Panel Content Protocol schemas ───────────────────────────────────────


class TestPanelContentProtocol:
    def test_panel_type_enum(self):
        assert PanelType.file_preview.value == "file_preview"
        assert PanelType.task_output.value == "task_output"
        assert PanelType.html_app.value == "html_app"

    def test_panel_content_type_enum(self):
        assert PanelContentType.structured.value == "structured"
        assert PanelContentType.html_sandbox.value == "html_sandbox"

    def test_structured_widget_enum(self):
        assert StructuredWidget.table.value == "table"
        assert StructuredWidget.markdown.value == "markdown"
        assert StructuredWidget.image.value == "image"
        assert StructuredWidget.json.value == "json"

    def test_structured_content_model(self):
        sc = StructuredContent(
            widget="table",
            data={
                "headers": ["Name", "Age"],
                "rows": [["Alice", "30"]],
            },
        )
        assert sc.widget == "table"
        assert len(sc.data["rows"]) == 1

    def test_html_sandbox_content_defaults(self):
        hsc = HtmlSandboxContent(html="<div>Hello</div>")
        assert hsc.height == 400
        assert hsc.permissions == []

    def test_panel_content_structured(self):
        pc = PanelContent(
            version="1",
            type="structured",
            title="Analysis Result",
            structured=StructuredContent(
                widget="table",
                data={"headers": ["Col"], "rows": [["A"]]},
            ),
        )
        d = pc.model_dump()
        assert d["version"] == "1"
        assert d["type"] == "structured"
        assert d["structured"]["widget"] == "table"
        assert d["html_sandbox"] is None

    def test_panel_content_html_sandbox_schema(self):
        """html_sandbox schema can be constructed even though not implemented."""
        pc = PanelContent(
            version="1",
            type="html_sandbox",
            title="Skill UI",
            html_sandbox=HtmlSandboxContent(
                html="<div>Custom</div>",
                height=600,
                permissions=["clipboard"],
            ),
        )
        d = pc.model_dump()
        assert d["type"] == "html_sandbox"
        assert d["html_sandbox"]["height"] == 600
        assert d["structured"] is None

    def test_panel_content_defaults(self):
        pc = PanelContent(title="Test")
        assert pc.version == "1"
        assert pc.type == "structured"
        assert pc.subtitle is None

    def test_panel_update_event_shape(self):
        """Validate the shape of a panel_update SSE event."""
        event = {
            "event": "panel_update",
            "panel_type": "file_preview",
            "panel_content": PanelContent(
                title="report.docx",
                structured=StructuredContent(
                    widget="json",
                    data={"kind": "docx", "paragraph_count": 10},
                ),
            ).model_dump(),
        }
        assert event["event"] == "panel_update"
        assert event["panel_type"] == "file_preview"
        assert event["panel_content"]["title"] == "report.docx"
