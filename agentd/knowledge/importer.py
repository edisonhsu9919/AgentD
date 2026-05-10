"""System-level knowledge import (Phase 6E — restructured).

Two-stage import workflow:
  Stage 1 (API call): User confirms metadata → raw file copied → background task starts
  Stage 2 (background): Extract content → generate Markdown → commit to knowledge store

Progress is persisted to .agentd/knowledge_imports/{task_id}.json so it
survives panel close, logout, and session deletion.

No iframe. No skill dependency. No model command assembly.
"""

import asyncio
import json
import logging
import os
import shutil
import uuid
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)

# ── Progress persistence ────────────────────────────────────────────────────

IMPORTS_DIR = ".agentd/knowledge_imports"


def _get_progress_path(session_dir: str, task_id: str) -> str:
    dir_path = os.path.join(session_dir, IMPORTS_DIR)
    os.makedirs(dir_path, exist_ok=True)
    return os.path.join(dir_path, f"{task_id}.json")


def _write_progress(session_dir: str, task_id: str, progress: dict) -> None:
    progress["updated_at"] = datetime.now(timezone.utc).isoformat()
    path = _get_progress_path(session_dir, task_id)
    with open(path, "w", encoding="utf-8") as f:
        json.dump(progress, f, ensure_ascii=False, indent=2)


def read_import_progress(session_dir: str, task_id: str) -> dict | None:
    path = _get_progress_path(session_dir, task_id)
    if not os.path.isfile(path):
        return None
    with open(path, "r", encoding="utf-8") as f:
        return json.load(f)


def list_import_tasks(session_dir: str) -> list[dict]:
    dir_path = os.path.join(session_dir, IMPORTS_DIR)
    if not os.path.isdir(dir_path):
        return []
    results = []
    for fname in sorted(os.listdir(dir_path)):
        if not fname.endswith(".json"):
            continue
        path = os.path.join(dir_path, fname)
        try:
            with open(path, "r", encoding="utf-8") as f:
                results.append(json.load(f))
        except Exception:
            pass
    results.sort(key=lambda x: x.get("created_at", ""), reverse=True)
    return results


# ── File extraction ─────────────────────────────────────────────────────────

def extract_content(file_path: str) -> str:
    ext = Path(file_path).suffix.lower()
    if ext == ".pdf":
        return _extract_pdf(file_path)
    elif ext == ".docx":
        return _extract_docx(file_path)
    elif ext == ".pptx":
        return _extract_pptx(file_path)
    elif ext in (".txt", ".md"):
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            return f.read()
    elif ext in (".png", ".jpg", ".jpeg", ".webp", ".bmp"):
        return f"(Image file: {Path(file_path).name}. OCR extraction requires VLM.)"
    return f"(Unsupported file type: {ext})"


def _extract_pdf(path: str) -> str:
    from pypdf import PdfReader
    reader = PdfReader(path)
    parts = []
    for i, page in enumerate(reader.pages):
        try:
            text = (page.extract_text() or "").strip()
        except Exception:
            text = ""
        if text:
            parts.append(f"## Page {i + 1}\n\n{text}")
    return "\n\n".join(parts) if parts else "(PDF contained no extractable text)"


def _extract_docx(path: str) -> str:
    from docx import Document
    doc = Document(path)
    parts = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        if para.style and para.style.name.startswith("Heading"):
            level = para.style.name.replace("Heading", "").strip()
            try:
                level = int(level)
            except ValueError:
                level = 2
            parts.append(f"{'#' * (level + 1)} {text}")
        else:
            parts.append(text)
    return "\n\n".join(parts) if parts else "(DOCX contained no text)"


def _extract_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_parts = [f"## Slide {i}"]
        if slide.shapes.title and slide.shapes.title.text:
            slide_parts.append(f"### {slide.shapes.title.text}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text and text != (slide.shapes.title.text if slide.shapes.title else ""):
                        slide_parts.append(text)
        if len(slide_parts) > 1:
            parts.append("\n\n".join(slide_parts))
    return "\n\n".join(parts) if parts else "(PPTX contained no text)"


# ── Metadata draft ──────────────────────────────────────────────────────────

DESCRIPTION_MAX_CHARS = 200

_DRAFT_PROMPT = """\
Based on the following document inspection result, generate metadata for a knowledge base entry.

Respond with ONLY a valid JSON object containing these fields:
- "title": A clear, concise document title (max 100 chars)
- "description": A brief summary of what this document is about (120-200 chars)
- "tags": An array of 3-5 keyword tags relevant to the content

Document info:
- Filename: {filename}
- Type: {kind}
- Size: {file_size} bytes

Content sample:
{content_sample}

Output ONLY the JSON object, no markdown fences, no explanation."""


async def generate_metadata_draft(file_path: str) -> dict[str, Any]:
    """Generate metadata suggestions using file_inspect + LLM.

    1. Uses the workspace/inspect extraction layer for a lightweight content sample
    2. Sends the sample to VLM/LLM for structured title/description/tags extraction
    3. Falls back to filename-based draft if LLM unavailable
    """
    filename = os.path.basename(file_path)
    stem = Path(filename).stem.replace("_", " ").replace("-", " ")
    ext = Path(filename).suffix.lower()
    kind_map = {".pdf": "pdf", ".docx": "docx", ".pptx": "pptx",
                ".png": "image", ".jpg": "image", ".jpeg": "image",
                ".txt": "text", ".md": "text"}
    kind = kind_map.get(ext, "unknown")
    file_size = os.path.getsize(file_path) if os.path.isfile(file_path) else 0

    # Step 1: Get content sample via file_inspect extraction layer
    content_sample = _get_content_sample(file_path, ext)

    # Step 2: Try LLM extraction
    llm_draft = await _llm_extract_metadata(filename, kind, file_size, content_sample)

    if llm_draft:
        return {
            "title": llm_draft.get("title", stem)[:100],
            "description": llm_draft.get("description", f"Knowledge document: {stem}")[:DESCRIPTION_MAX_CHARS],
            "tags": llm_draft.get("tags", [])[:5],
            "permission": "private",
            "kind": kind,
            "filename": filename,
            "file_size": file_size,
            "limits": {"description_max_chars": DESCRIPTION_MAX_CHARS},
        }

    # Fallback: filename-based
    return {
        "title": stem[:100],
        "description": f"Knowledge document: {stem}"[:DESCRIPTION_MAX_CHARS],
        "tags": [],
        "permission": "private",
        "kind": kind,
        "filename": filename,
        "file_size": file_size,
        "limits": {"description_max_chars": DESCRIPTION_MAX_CHARS},
    }


def _get_content_sample(file_path: str, ext: str) -> str:
    """Get a lightweight content sample using the file_inspect extraction layer."""
    try:
        from workspace.router import _extract_file_info
        import asyncio

        # _extract_file_info is async but we need sync here for the draft
        # Use the sync extraction functions directly
        if ext == ".pdf":
            from files.pdf import extract
            result = extract(file_path)
            return result.get("text_sample", "")[:2000]
        elif ext == ".docx":
            from files.office_docx import extract
            result = extract(file_path)
            return result.get("text_sample", "")[:2000]
        elif ext == ".pptx":
            from files.office_pptx import extract
            result = extract(file_path)
            slides = result.get("slides", [])
            parts = []
            for s in slides[:5]:
                if s.get("title"):
                    parts.append(s["title"])
                for b in s.get("bullets", [])[:3]:
                    parts.append(f"  - {b}")
            return "\n".join(parts)[:2000]
        elif ext in (".txt", ".md"):
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                return f.read(2000)
        elif ext == ".eml":
            from files.email_eml import extract
            result = extract(file_path)
            return f"Subject: {result.get('subject', '')}\n{result.get('body_preview', '')}"[:2000]
    except Exception as e:
        logger.warning("Content sample extraction failed: %s", e)

    return ""


async def _llm_extract_metadata(
    filename: str, kind: str, file_size: int, content_sample: str,
) -> dict | None:
    """Call VLM/LLM to extract structured metadata from content sample."""
    if not content_sample.strip():
        return None

    try:
        from core.database import AsyncSessionLocal
        from model_config.service import resolve_active_vlm_config
        import httpx

        async with AsyncSessionLocal() as db:
            resolved = await resolve_active_vlm_config(db)

        if resolved is None:
            # No VLM — try main LLM
            from model_config.service import resolve_active_model_config
            async with AsyncSessionLocal() as db:
                resolved = await resolve_active_model_config(db)

        prompt = _DRAFT_PROMPT.format(
            filename=filename,
            kind=kind,
            file_size=file_size,
            content_sample=content_sample[:2000],
        )

        payload = {
            "model": resolved.model_id,
            "messages": [{"role": "user", "content": prompt}],
            "max_tokens": 512,
        }
        headers = {"Content-Type": "application/json"}
        if resolved.api_key:
            headers["Authorization"] = f"Bearer {resolved.api_key}"

        url = resolved.base_url.rstrip("/") + "/chat/completions"

        async with httpx.AsyncClient(trust_env=False, timeout=30.0) as client:
            resp = await client.post(url, json=payload, headers=headers)

        if resp.status_code != 200:
            return None

        body = resp.json()
        choices = body.get("choices", [])
        if not choices:
            return None

        text = choices[0].get("message", {}).get("content", "")
        # Try to parse JSON from response
        text = text.strip()
        if text.startswith("```"):
            text = text.split("\n", 1)[-1].rsplit("```", 1)[0].strip()

        return json.loads(text)

    except Exception as e:
        logger.warning("LLM metadata extraction failed: %s", e)
        return None


# ── Background import task ──────────────────────────────────────────────────

async def run_import_task(
    *,
    task_id: str,
    session_id: str,
    session_dir: str,
    user_id: str,
    source_path: str,
    raw_path: str,
    metadata: dict,
    publish_fn=None,
) -> dict[str, Any]:
    """Background task: extract content → build Markdown → commit.

    Raw file is already copied before this runs.
    Progress is persisted so frontend can poll.
    """
    filename = os.path.basename(source_path)
    kind = metadata.get("kind", "unknown")
    title = metadata.get("title", filename)
    created = datetime.now(timezone.utc).isoformat()

    # Stable progress template — all fields always present
    def _progress(status: str, phase: str, **extra) -> dict:
        base = {
            "task_id": task_id,
            "status": status,
            "phase": phase,
            "filename": filename,
            "kind": kind,
            "title": title,
            "source_path": os.path.basename(source_path),
            "raw_path": f"knowledge/raw/{filename}",
            "content_chars": 0,
            "doc_id": None,
            "error": None,
            "created_at": created,
        }
        base.update(extra)
        return base

    try:
        # Phase 1: Extracting
        _write_progress(session_dir, task_id, _progress("extracting", "Extracting content..."))
        print(f"[knowledge-import] Extracting: {filename}")

        content = extract_content(raw_path)

        # Phase 2: Committing
        _write_progress(session_dir, task_id, _progress(
            "committing", "Writing to knowledge base...", content_chars=len(content),
        ))
        print(f"[knowledge-import] Committing: {filename} ({len(content)} chars)")

        from knowledge.store import (
            build_frontmatter,
            ensure_knowledge_dirs,
            generate_doc_id,
            validate_frontmatter,
            write_knowledge_doc,
        )

        ensure_knowledge_dirs()
        doc_id = generate_doc_id()

        tags = metadata.get("tags", [])
        if isinstance(tags, str):
            tags = [t.strip() for t in tags.split(",") if t.strip()]

        fm = build_frontmatter(
            title=metadata.get("title", filename),
            description=metadata.get("description", ""),
            kind=kind,
            owner=user_id,
            permission=metadata.get("permission", "private"),
            tags=tags,
            source_file=filename,
            source_path=f"knowledge/raw/{filename}",
            file_size=os.path.getsize(raw_path) if os.path.isfile(raw_path) else 0,
        )
        errors = validate_frontmatter(fm)
        if errors:
            raise ValueError("; ".join(errors))

        write_knowledge_doc(doc_id, fm, content)

        # Phase 3: Done
        _write_progress(session_dir, task_id, _progress(
            "completed", "Import complete",
            content_chars=len(content), doc_id=doc_id, title=fm["title"],
        ))
        print(f"[knowledge-import] Done: doc_id={doc_id}")

        if publish_fn:
            try:
                await publish_fn(session_id, {
                    "event": "knowledge_import_done",
                    "task_id": task_id,
                    "doc_id": doc_id,
                    "title": fm["title"],
                })
            except Exception:
                pass

        return {"success": True, "doc_id": doc_id, "title": fm["title"]}

    except Exception as e:
        logger.error("Import task %s failed: %s", task_id[:8], e, exc_info=True)
        _write_progress(session_dir, task_id, _progress(
            "failed", f"Error: {e}", error=str(e),
        ))
        return {"success": False, "reason": str(e)}
