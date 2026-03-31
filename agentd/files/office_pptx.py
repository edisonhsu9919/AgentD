"""PPTX extraction layer (Phase O2).

Pure data extraction from PPTX files using python-pptx.
No LLM calls — structured raw data only.
"""

import os
from typing import Any

# Limits
_SLIDE_TEXT_MAX_CHARS = 300
_MAX_SLIDES_DETAIL = 30


def extract(path: str) -> dict[str, Any]:
    """Extract structural metadata and slide summaries from a PPTX file.

    Returns a dict with:
      - path, kind, office_kind, size_bytes
      - slide_count
      - slides (list of per-slide info: number, title, text_preview, has_notes)
      - metadata (title, author, subject)

    Raises FileNotFoundError if path doesn't exist.
    Raises ValueError if file is not a valid PPTX.
    """
    if not os.path.isfile(path):
        raise FileNotFoundError(f"File not found: {path}")

    size_bytes = os.path.getsize(path)

    try:
        from pptx import Presentation
    except ImportError as e:
        raise RuntimeError("python-pptx is not installed. Run: pip install python-pptx") from e

    try:
        prs = Presentation(path)
    except Exception as e:
        raise ValueError(f"Cannot read PPTX: {e}") from e

    # Metadata
    props = prs.core_properties
    metadata = {
        "title": _clean(props.title),
        "author": _clean(props.author),
        "subject": _clean(props.subject),
    }

    # Slides
    slide_count = len(prs.slides)
    slides_info: list[dict[str, Any]] = []

    for i, slide in enumerate(prs.slides):
        if i >= _MAX_SLIDES_DETAIL:
            break
        slides_info.append(_extract_slide(slide, i + 1))

    return {
        "path": path,
        "kind": "office",
        "office_kind": "pptx",
        "size_bytes": size_bytes,
        "slide_count": slide_count,
        "slides": slides_info,
        "metadata": metadata,
    }


def _extract_slide(slide, number: int) -> dict[str, Any]:
    """Extract info from a single slide."""
    title = ""
    texts: list[str] = []

    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if shape.shape_id == slide.shapes.title.shape_id if slide.shapes.title else False:
                title = text
            else:
                texts.append(text)

    # If title wasn't found via shape_id matching, try the title placeholder
    if not title and slide.shapes.title:
        title = slide.shapes.title.text.strip()

    # Text preview
    full_text = "\n".join(texts)
    text_preview = full_text[:_SLIDE_TEXT_MAX_CHARS]

    # Notes
    has_notes = False
    if slide.has_notes_slide:
        notes_text = slide.notes_slide.notes_text_frame.text.strip()
        has_notes = len(notes_text) > 0

    return {
        "number": number,
        "title": title,
        "text_preview": text_preview,
        "has_notes": has_notes,
    }


def _clean(value) -> str:
    if value is None:
        return ""
    return str(value).strip()
